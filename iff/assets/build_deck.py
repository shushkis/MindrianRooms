#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
IRIS Week 2 Deck Builder — IFF + Sensor Fusion Platform
Eli + Mordi | 6 slides | Hebrew
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
BG_CARD   = RGBColor(0x1E, 0x2A, 0x45)   # card navy
ACCENT1   = RGBColor(0x38, 0x8B, 0xFD)   # blue
ACCENT2   = RGBColor(0x10, 0xB9, 0x81)   # teal/green
ACCENT3   = RGBColor(0xF5, 0x9E, 0x0B)   # amber
RED_ACC   = RGBColor(0xEF, 0x44, 0x44)   # red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LT   = RGBColor(0x94, 0xA3, 0xB8)
GRAY_MED  = RGBColor(0x64, 0x74, 0x8B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def make_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, opacity=1.0):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def set_rtl(para):
    """Make a paragraph RTL."""
    pPr = para._pPr
    if pPr is None:
        pPr = para._p.get_or_add_pPr()
    pPr.set(qn('a:rtl'), '1')

def add_textbox(slide, l, t, w, h, text, font_size, color,
                bold=False, align=PP_ALIGN.RIGHT, rtl=True,
                italic=False, line_spacing=None):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if rtl:
        set_rtl(p)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    if line_spacing:
        p.line_spacing = line_spacing
    return txBox

def add_multiline(slide, l, t, w, h, lines, font_size, color,
                  bold=False, align=PP_ALIGN.RIGHT, rtl=True,
                  leading_pt=None, first_bold=False):
    """Add a textbox with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if rtl:
            set_rtl(p)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        is_bold = bold or (first_bold and i == 0)
        run.font.bold = is_bold
        if leading_pt:
            p.space_before = Pt(leading_pt)
    return txBox

def accent_bar(slide, color, x=Inches(0.4), y=Inches(0.35), w=Inches(0.08), h=Inches(0.9)):
    add_rect(slide, x, y, w, h, color)

def slide_number(slide, n, total=6):
    add_textbox(slide,
                SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45),
                Inches(1.0), Inches(0.35),
                f"{n}/{total}", 10, GRAY_MED, align=PP_ALIGN.RIGHT, rtl=False)

def divider(slide, y, color=GRAY_MED):
    line = slide.shapes.add_connector(1, Inches(0.4), y, SLIDE_W - Inches(0.4), y)
    line.line.color.rgb = color
    line.line.width = Pt(0.5)

# ── Slide builders ────────────────────────────────────────────────────────────

def slide1_opening(prs):
    sl = blank_slide(prs)
    fill_bg(sl, BG_DARK)
    accent_bar(sl, ACCENT1)

    # Top label
    add_textbox(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.4),
                "IRIS שבוע 2  ·  שינוי בלב החדשנות  ·  אלי + מורדי",
                10, GRAY_LT, align=PP_ALIGN.RIGHT)

    # Main title
    add_textbox(sl, Inches(0.6), Inches(1.1), Inches(12.3), Inches(0.9),
                "ליבה טכנולוגית אחת — שלושה שווקים שבהם הבעיה לא נפתרה",
                30, WHITE, bold=True, align=PP_ALIGN.RIGHT)

    # Sub
    add_textbox(sl, Inches(0.6), Inches(1.95), Inches(12.3), Inches(0.55),
                "Sensor Fusion  ·  Computer Vision  ·  GPS-Denied  ·  Unstructured Environments",
                14, ACCENT1, align=PP_ALIGN.RIGHT, rtl=False)

    divider(sl, Inches(2.7), ACCENT1)

    # Two columns: Eli | Mordi
    # Eli block
    add_rect(sl, Inches(7.0), Inches(2.9), Inches(5.8), Inches(1.8), BG_CARD)
    add_textbox(sl, Inches(7.1), Inches(2.95), Inches(5.6), Inches(0.4),
                "אלי", 16, ACCENT1, bold=True)
    add_multiline(sl, Inches(7.1), Inches(3.35), Inches(5.6), Inches(1.3),
                  ["11 שנים — Sensor Fusion, Computer Vision, AV",
                   "מחלקת מכשור  ·  טייסת ניסוי טיסה"],
                  12, GRAY_LT)

    # Mordi block
    add_rect(sl, Inches(0.5), Inches(2.9), Inches(6.3), Inches(1.8), BG_CARD)
    add_textbox(sl, Inches(0.6), Inches(2.95), Inches(6.1), Inches(0.4),
                "מורדי", 16, ACCENT2, bold=True)
    add_multiline(sl, Inches(0.6), Inches(3.35), Inches(6.1), Inches(1.3),
                  ["מ\"פ צנחנים  ·  3 תיאטרות (עזה, לבנון, סוריה)",
                   "RZR בשטח קשה  ·  עורך דין, רישום מקרקעין"],
                  12, GRAY_LT)

    # Insight box
    add_rect(sl, Inches(0.5), Inches(5.05), Inches(12.3), Inches(1.1), BG_CARD)
    add_textbox(sl, Inches(0.6), Inches(5.1), Inches(12.0), Inches(0.45),
                "המכנה המשותף לשלוש ההזדמנויות:", 13, ACCENT1, bold=True)
    add_textbox(sl, Inches(0.6), Inches(5.5), Inches(12.0), Inches(0.55),
                "תנאים שבהם כלים קיימים מניחים שמשאבים קיימים — GPS, תשתית, ראות, זמן — ובשטח הם לא.",
                13, WHITE)

    # Footer note
    add_textbox(sl, Inches(0.5), Inches(6.3), Inches(9), Inches(0.4),
                "מי שבנה Sensor Fusion לרכבים אוטונומיים — כבר פתר 60% מהאתגר הטכני בשלוש ההזדמנויות האלה.",
                11, GRAY_LT)
    slide_number(sl, 1)
    return sl


def slide2_opp1(prs):
    sl = blank_slide(prs)
    fill_bg(sl, BG_DARK)
    accent_bar(sl, ACCENT1)

    add_textbox(sl, Inches(0.6), Inches(0.25), Inches(4.5), Inches(0.4),
                "הזדמנות 1", 11, ACCENT1, bold=True)
    add_textbox(sl, Inches(0.6), Inches(0.65), Inches(12.3), Inches(0.75),
                "IFF בנקודת הירי — AI + UWB על הנשק",
                26, WHITE, bold=True)
    add_textbox(sl, Inches(0.6), Inches(1.35), Inches(12.3), Inches(0.4),
                "SmarTrack מגיד לך איפה החברים. לא מגיד לך אם מי שבכוונת — חבר.",
                13, GRAY_LT, italic=True)

    divider(sl, Inches(1.85))

    # Left: Data column
    add_textbox(sl, Inches(7.5), Inches(1.95), Inches(5.4), Inches(0.4),
                "הנתונים", 12, ACCENT1, bold=True)
    data_lines = [
        "15–20% מנפגעי מלחמות העולם + וייטנאם = אש ידידותית",
        "Desert Storm: 24% מהרוגי קרב = אש ידידותית",
        "Army: ≤9 אינץ\", ≥328 יארד, ≥270°, עמיד NVG",
        "אין מוצר שעומד בכל הדרישות יחד (2025)",
        "TurbineOne: $99M Army AI — מ-ISR, לא מנשק"
    ]
    add_multiline(sl, Inches(7.5), Inches(2.35), Inches(5.4), Inches(2.0),
                  data_lines, 10.5, GRAY_LT, leading_pt=3)

    # Right: Two mechanism cards
    add_rect(sl, Inches(0.5), Inches(1.95), Inches(6.7), Inches(0.9), BG_CARD)
    add_textbox(sl, Inches(0.6), Inches(2.0), Inches(3.0), Inches(0.35),
                "מנגנון 1 — AI Visual", 11, ACCENT1, bold=True)
    add_textbox(sl, Inches(0.6), Inches(2.3), Inches(6.5), Inches(0.5),
                "Computer Vision + Thermal Fusion  ·  מזהה חתימה ידידותית  ·  ללא emission  ·  edge-only",
                10, GRAY_LT, rtl=False)

    add_rect(sl, Inches(0.5), Inches(3.0), Inches(6.7), Inches(0.9), BG_CARD)
    add_textbox(sl, Inches(0.6), Inches(3.05), Inches(3.5), Inches(0.35),
                "מנגנון 2 — UWB Directional", 11, ACCENT2, bold=True)
    add_textbox(sl, Inches(0.6), Inches(3.35), Inches(6.5), Inches(0.5),
                "IEEE 802.15.4z  ·  דיוק סנטימטרי  ·  GPS-Denied  ·  \"מטרה בכוונת = חבר ב-network\"",
                10, GRAY_LT, rtl=False)

    add_textbox(sl, Inches(0.5), Inches(4.05), Inches(6.7), Inches(0.35),
                "CV עובד ללא network  ·  UWB עובד בחשכה מוחלטת  →  יחד: כיסוי מלא",
                10.5, ACCENT2, italic=True)

    # WD formulations
    divider(sl, Inches(4.5))
    add_textbox(sl, Inches(0.5), Inches(4.55), Inches(3), Inches(0.35),
                "ניסוחי בעיה", 11, ACCENT1, bold=True)

    wd_lines = [
        "WD.1 | מ\"פ בבניין, לילה, עשן  →  מדד: time-to-IFF < 100ms",
        "WD.2 | RZR, GNSS jamming, כניסה מ-2 צירים  →  מדד: autonomous op time ללא GNSS",
        "WD.3 | מחסום, רכב מתקרב, 2 שניות  →  מדד: false positive + false negative rate"
    ]
    add_multiline(sl, Inches(0.5), Inches(4.95), Inches(12.3), Inches(1.4),
                  wd_lines, 10.5, GRAY_LT, leading_pt=4)

    slide_number(sl, 2)
    return sl


def slide3_opp2(prs):
    sl = blank_slide(prs)
    fill_bg(sl, BG_DARK)
    accent_bar(sl, ACCENT2)

    add_textbox(sl, Inches(0.6), Inches(0.25), Inches(4.5), Inches(0.4),
                "הזדמנות 2", 11, ACCENT2, bold=True)
    add_textbox(sl, Inches(0.6), Inches(0.65), Inches(12.3), Inches(0.75),
                "ראיות תצלומי אויר AI — מנוע לבית משפט",
                26, WHITE, bold=True)
    add_textbox(sl, Inches(0.6), Inches(1.35), Inches(12.3), Inches(0.4),
                "AI מנתח נכסים לצרכי מיסוי מזה שנים. לצרכי בית משפט — עדיין שבועות של עבודה ידנית.",
                13, GRAY_LT, italic=True)

    divider(sl, Inches(1.85))

    # Data
    add_textbox(sl, Inches(7.5), Inches(1.95), Inches(5.4), Inches(0.4),
                "הנתונים", 12, ACCENT2, bold=True)
    data_lines = [
        "AI change-detection נרחב לצרכי מיסוי (Nearmap, Vexcel, FlyPix)",
        "GEOward — 23 שנות ניסיון — לא מפרסמת מתודולוגיה, לא QA",
        "\"AI-enhanced evidence — legal systems are struggling\" (Opinio Juris, 2025)",
        "קושן טורקי + מאגר השמיים הלאומי = dataset ייחודי לישראל"
    ]
    add_multiline(sl, Inches(7.5), Inches(2.35), Inches(5.4), Inches(1.8),
                  data_lines, 10.5, GRAY_LT, leading_pt=4)

    # Advantage box
    add_rect(sl, Inches(7.5), Inches(4.2), Inches(5.4), Inches(0.9), BG_CARD)
    add_textbox(sl, Inches(7.6), Inches(4.25), Inches(5.2), Inches(0.35),
                "Unfair Advantage", 11, ACCENT2, bold=True)
    add_textbox(sl, Inches(7.6), Inches(4.55), Inches(5.2), Inches(0.45),
                "מורדי = היחיד בשוק שהוא גם מומחה ראייתי וגם מכיר מה שופט מקבל ומה הוא פוסל.",
                10.5, WHITE)

    # S-curve note
    add_rect(sl, Inches(0.5), Inches(1.95), Inches(6.7), Inches(1.5), BG_CARD)
    add_textbox(sl, Inches(0.6), Inches(2.0), Inches(6.5), Inches(0.4),
                "עקומת S — הנישה שלא עלתה עדיין", 12, ACCENT2, bold=True)
    add_multiline(sl, Inches(0.6), Inches(2.4), Inches(6.5), Inches(1.0),
                  ["שלב א': ניתוח ידני בלבד — שבועות לכל תיק",
                   "שלב ב': AI לצרכי מיסוי — נרחב ומוכח",
                   "הפער: הנישה הראייתית/משפטית — עדיין ידנית ב-2025 ↑"],
                  11, GRAY_LT, leading_pt=3)

    # WD
    divider(sl, Inches(3.6))
    add_textbox(sl, Inches(0.5), Inches(3.65), Inches(3), Inches(0.35),
                "ניסוחי בעיה", 11, ACCENT2, bold=True)

    wd_lines = [
        "WD.1 | עורך דין, מחלוקת בעלות חקלאית, קושן טורקי  →  מדד: דוח תוך 72 שעות, עלות <₪5,000, עמיד בחקירה",
        "WD.2 | מומחה GIS כעד מומחה, מאגר תמונות לא ממוין  →  מדד: שיעור ניצחון בביהמ\"ש",
        "WD.3 | שמאי + עיכוב סגירת עסקה, קושן ללא גיאורפרנסינג  →  מדד: ימי עיכוב בסגירה"
    ]
    add_multiline(sl, Inches(0.5), Inches(4.05), Inches(12.3), Inches(2.0),
                  wd_lines, 10.5, GRAY_LT, leading_pt=5)

    slide_number(sl, 3)
    return sl


def slide4_opp3(prs):
    sl = blank_slide(prs)
    fill_bg(sl, BG_DARK)
    accent_bar(sl, ACCENT3)

    add_textbox(sl, Inches(0.6), Inches(0.25), Inches(4.5), Inches(0.4),
                "הזדמנות 3", 11, ACCENT3, bold=True)
    add_textbox(sl, Inches(0.6), Inches(0.65), Inches(12.3), Inches(0.75),
                "AV לשטח לא מובנה — מיפוי שטח צבאי וחקלאי",
                26, WHITE, bold=True)
    add_textbox(sl, Inches(0.6), Inches(1.35), Inches(12.3), Inches(0.4),
                "5–10 שנים מאחורי ה-AV הרגיל — ועדיין לא מישהו שמוכר לחקלאי, לכורה, לצבא.",
                13, GRAY_LT, italic=True)

    divider(sl, Inches(1.85))

    # Data
    add_textbox(sl, Inches(7.5), Inches(1.95), Inches(5.4), Inches(0.4),
                "הנתונים", 12, ACCENT3, bold=True)
    data_lines = [
        "Army: $15.5M ל-3 startups לנסות off-road AV (ספטמבר 2025)",
        "\"5–10 years behind urban AV\" (Wiley JFR, 2025)",
        "Solid-state LiDAR: $500–$10,000 ב-2025 — עלות כבר לא חסם",
        "95%+ ממחקר ה-AV: כבישים מסומנים בלבד"
    ]
    add_multiline(sl, Inches(7.5), Inches(2.35), Inches(5.4), Inches(1.8),
                  data_lines, 10.5, GRAY_LT, leading_pt=4)

    # Connection box
    add_rect(sl, Inches(7.5), Inches(4.3), Inches(5.4), Inches(0.85), BG_CARD)
    add_textbox(sl, Inches(7.6), Inches(4.35), Inches(5.2), Inches(0.35),
                "החיבור לטכנולוגיית IFF", 11, ACCENT3, bold=True)
    add_textbox(sl, Inches(7.6), Inches(4.65), Inches(5.2), Inches(0.45),
                "אותה ארכיטקטורה: Sensor Fusion ללא GPS + Edge Inference + GPS-Denied Localization.",
                10.5, WHITE)

    # Three markets
    add_rect(sl, Inches(0.5), Inches(1.95), Inches(6.7), Inches(1.5), BG_CARD)
    add_textbox(sl, Inches(0.6), Inches(2.0), Inches(6.5), Inches(0.4),
                "שלושה שווקים, אותה טכנולוגיה", 12, ACCENT3, bold=True)
    add_multiline(sl, Inches(0.6), Inches(2.4), Inches(6.5), Inches(1.0),
                  ["צבאי — RZR, UGV, סיור בשטח בלתי-מסומן",
                   "חקלאי — סיור גבולות, ריסוס, מעקב",
                   "כרייה — כלים כבדים בתנאי אבק, אפס GPS"],
                  11, GRAY_LT, leading_pt=3)

    # WD
    divider(sl, Inches(3.6))
    add_textbox(sl, Inches(0.5), Inches(3.65), Inches(3), Inches(0.35),
                "ניסוחי בעיה", 11, ACCENT3, bold=True)

    wd_lines = [
        "WD.1 | חקלאי — 500 דונם, שטח פתוח, אין תשתית  →  מדד: שעות אוטונומיה ביום vs. עלות שמירה ידנית",
        "WD.2 | גוף ביטחוני — RZR, GNSS jamming, חשכה  →  מדד: autonomous op time ללא GNSS",
        "WD.3 | מכרה פתוח — אבק כבד, LiDAR מתעוור  →  מדד: שעות אוטונומיות למשמרת (בסיס: <2 שעות)"
    ]
    add_multiline(sl, Inches(0.5), Inches(4.05), Inches(12.3), Inches(2.0),
                  wd_lines, 10.5, GRAY_LT, leading_pt=5)

    slide_number(sl, 4)
    return sl


def slide5_feedback(prs):
    sl = blank_slide(prs)
    fill_bg(sl, BG_DARK)
    accent_bar(sl, WHITE)

    add_textbox(sl, Inches(0.6), Inches(0.25), Inches(12.3), Inches(0.5),
                "מה אנחנו רוצים מהקוהורט", 11, GRAY_LT)
    add_textbox(sl, Inches(0.6), Inches(0.65), Inches(12.3), Inches(0.65),
                "3 שאלות שרק החדר הזה יכול לענות עליהן",
                26, WHITE, bold=True)

    divider(sl, Inches(1.5))

    cards = [
        (ACCENT1, "שאלה 1 — על הבחירה",
         "מי בחדר ראה צוות עם \"פלטפורמה רוחבית\" שניסה לעבוד בכמה שווקים בו-זמנית — ומה קרה?"),
        (ACCENT2, "שאלה 2 — על ה-Beachhead",
         "מי כבר סגר חוזה POC עם לקוח ביטחוני ישראלי בפחות מ-6 חודשים — ומה היה המסלול?\n"
         "(IFF: מחזור רכש ארוך  ·  ראיות אויר: עורך דין ב-90 יום  ·  AV שטח: שניים שונים מאוד)"),
        (ACCENT3, "שאלה 3 — על ה-Moat",
         "Elbit, Rafael, IAI — 18 חודשים לחיקוי. Nearmap — גם כן.\n"
         "מהי ההגנה האמיתית: הטכנולוגיה, נתוני האימון, גישה לשטח, או המוניטין המבצעי?"),
    ]

    y_positions = [Inches(1.6), Inches(3.0), Inches(4.65)]
    for i, (color, title, body) in enumerate(cards):
        y = y_positions[i]
        add_rect(sl, Inches(0.5), y, Inches(12.3), Inches(1.2), BG_CARD)
        add_rect(sl, Inches(0.5), y, Inches(0.07), Inches(1.2), color)
        add_textbox(sl, Inches(0.7), y + Pt(6), Inches(12.0), Inches(0.38),
                    title, 13, color, bold=True)
        add_textbox(sl, Inches(0.7), y + Inches(0.42), Inches(12.0), Inches(0.72),
                    body, 11, GRAY_LT)

    slide_number(sl, 5)
    return sl


def slide6_doubts(prs):
    sl = blank_slide(prs)
    fill_bg(sl, BG_DARK)
    accent_bar(sl, RED_ACC)

    add_textbox(sl, Inches(0.6), Inches(0.25), Inches(12.3), Inches(0.4),
                "ספקות פתוחים", 11, RED_ACC, bold=True)
    add_textbox(sl, Inches(0.6), Inches(0.6), Inches(12.3), Inches(0.65),
                "מה יכול להרוג את כל זה — בגלוי",
                26, WHITE, bold=True)

    divider(sl, Inches(1.45))

    doubts = [
        (RED_ACC,
         "ספק 1 — האם \"פלטפורמה\" היא תירוץ לחוסר מיקוד?",
         "שלוש הזדמנויות = שלוש חברות שונות. \"אותה טכנולוגיה\" לא מגנה על כך.\n"
         "לא ידוע: איזו תקבל \"כן, הנה כסף\" ראשונה מלקוח אמיתי?"),
        (ACCENT3,
         "ספק 2 — Training Data הוא חסם אמיתי ב-IFF",
         "TurbineOne בנתה זאת על $36M Series B + Army contracts.\n"
         "לא ידוע: האם Synthetic Data מספיק לאמת ראשונית? ואם לא — מי מאשר גישה לנתוני אימון מסווגים?"),
        (ACCENT2,
         "ספק 3 — ראיות אויר: Feature או Venture?",
         "GEOward יכולה לשלב AI פנימית תוך 6–9 חודשים.\n"
         "לא ידוע: האם ניתן לרשום IP על מתודולוגיה ראייתית — או שזה רק know-how?"),
    ]

    y_positions = [Inches(1.55), Inches(3.0), Inches(4.6)]
    for i, (color, title, body) in enumerate(doubts):
        y = y_positions[i]
        add_rect(sl, Inches(0.5), y, Inches(12.3), Inches(1.25), BG_CARD)
        add_rect(sl, Inches(0.5), y, Inches(0.07), Inches(1.25), color)
        add_textbox(sl, Inches(0.7), y + Pt(6), Inches(12.0), Inches(0.38),
                    title, 13, color, bold=True)
        add_textbox(sl, Inches(0.7), y + Inches(0.44), Inches(12.0), Inches(0.72),
                    body, 11, GRAY_LT)

    slide_number(sl, 6)
    return sl


# ── Main ──────────────────────────────────────────────────────────────────────

def build():
    prs = make_prs()
    slide1_opening(prs)
    slide2_opp1(prs)
    slide3_opp2(prs)
    slide4_opp3(prs)
    slide5_feedback(prs)
    slide6_doubts(prs)

    out = r"C:\Users\eidlm\MindrianRooms\iff\assets\IRIS_Week2_Eli_Mordi.pptx"
    prs.save(out)
    print(f"Saved: {out}")

if __name__ == "__main__":
    build()
